import logging
import io

from PIL.Image import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from src.mmore.process.utils import clean_text, clean_image
from src.mmore.type import FileDescriptor
from .processor import Processor, ProcessorConfig, ProcessorResult
from typing import List, Dict, Any

logger = logging.getLogger(__name__)


class PPTXProcessor(Processor):
    """
    A processor for handling PPTX files. Extracts text, images, and notes from PowerPoint presentations.

    Attributes:
        files (List[FileDescriptor]): List of files to be processed.
        config (ProcessorConfig): Configuration for the processor.
    """

    def __init__(self, files, config=None):
        """
        Args:
            files (List[FileDescriptor]): List of files to process.
            config (ProcessorConfig, optional): Configuration for the processor. Defaults to None.
        """
        super().__init__(files, config=config or ProcessorConfig())

    def accepts(self, file: FileDescriptor) -> bool:
        """
        Args:
            file (FileDescriptor): The file descriptor to check.

        Returns:
            bool: True if the file is a PPTX file, False otherwise.
        """
        return file.file_extension.lower() in [".pptx"]

    def require_gpu(self) -> bool:
        """
        Returns:
            tuple: A tuple (False, False) indicating no GPU requirement for both standard and fast modes.
        """
        return False

    def process_one_file(self, file_path: str, fast: bool = False) -> ProcessorResult:
        """
        Process a single PPTX file. Extracts text, images, and notes from each slide.

        Args:
            file_path (str): Path to the PPTX file.

        Returns:
            dict: A dictionary containing processed text and images.

        The method processes each slide, extracting text and images from shapes,
        and extracts notes if present. The elements are sorted by their vertical position.
        """
        super().process_one_file(file_path, fast=fast)

        logger.info(f"Processing PowerPoint file: {file_path}")
        try:
            prs = Presentation(file_path)
        except Exception as e:
            logger.error(f"Failed to open PowerPoint file {file_path}: {e}")
            return self.create_sample([], [], file_path)

        all_text: list[str] = []
        embedded_images: list[Image.Image] = []

        try:
            for slide in prs.slides:
                # 1) Extract text and images from slide
                # Sort shapes by their vertical position
                shape_list = sorted(
                    (shape for shape in slide.shapes if hasattr(shape, "top")),
                    key=lambda s: s.top,
                )

                for shape in shape_list:
                    # Extract text from shape
                    if shape.has_text_frame:
                        cleaned_text = clean_text(shape.text)
                        if cleaned_text.strip():
                            all_text.append(cleaned_text)

                    # Extract images from shape
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            pil_image = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
                            if clean_image(pil_image):
                                embedded_images.append(pil_image)
                                all_text.append(self.config.attachment_tag)

                        except Exception as e:
                            logger.error(f"Failed to extract image from slide: {e}")

                # 2) Extract text from slide notes if present
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame
                    for paragraph in notes.paragraphs:
                        if paragraph.text:
                            cleaned = clean_text(paragraph.text)
                            if cleaned.strip():
                                all_text.append(cleaned)

        except Exception as e:
            logger.error(f"[PPTX] Error processing slides in {file_path}: {e}")

        return self.create_sample(all_text, embedded_images, file_path)
